import os
import zipfile
import time
from pptx import Presentation
import win32com.client
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import font, ttk
from PIL import Image, ImageTk
import sys

def resource_path(relative_path):
    """ Get the absolute path to the resource, works for dev and for PyInstaller """
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)

def close_powerpoint_instances():
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Quit()
    except Exception as e:
        print(f"Error closing PowerPoint instances: {e}")

def convert_ppsx_to_pptx(ppsx_file, pptx_file):
    try:
        close_powerpoint_instances()
        time.sleep(5)  # Give more time to ensure PowerPoint instances are closed
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        print(f"Opening file: {ppsx_file}")
        presentation = powerpoint.Presentations.Open(ppsx_file, WithWindow=False)
        print(f"Saving file as: {pptx_file}")
        presentation.SaveAs(pptx_file, 24)  # 24 oznacza format pptx
        presentation.Close()
        powerpoint.Quit()
        time.sleep(5)  # Additional delay to ensure file is saved properly
    except Exception as e:
        print(f"Error converting PPSX to PPTX: {e}")
        print(f"Error details: {e.args}")
        raise

def export_slides_to_images(pptx_file, output_folder, style="black", fullscreen=False):
    try:
        close_powerpoint_instances()
        time.sleep(5)

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        print(f"Opening file: {pptx_file}")
        presentation = powerpoint.Presentations.Open(pptx_file, WithWindow=False)

        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        for i, slide in enumerate(presentation.Slides):
            image_path = os.path.abspath(os.path.join(output_folder, f"slide_{i + 1}.png"))
            print(f"Exporting slide {i + 1} as {image_path}")
            slide.Export(image_path, "PNG")

        presentation.Close()
        powerpoint.Quit()
        time.sleep(5)  # Additional delay to ensure images are exported properly
    except Exception as e:
        print(f"Error exporting slides to images: {e}")
        raise

def convert_pptx_to_scorm(pptx_file, output_folder, title, author, use_slide_titles, use_mini_images, style="black", fullscreen=False):
    try:
        export_slides_to_images(pptx_file, output_folder, style, fullscreen)
        slide_images = [os.path.abspath(os.path.join(output_folder, f"slide_{i + 1}.png")) for i in range(len(Presentation(pptx_file).slides))]
        slide_titles = [slide.shapes.title.text if slide.shapes.title else f"Slide {i + 1}" for i, slide in enumerate(Presentation(pptx_file).slides)] if use_slide_titles else None
        generate_html_files(slide_images, output_folder, title, author, slide_titles, use_mini_images, style, fullscreen)
        scorm_manifest = generate_scorm_manifest(slide_images)
        manifest_path = os.path.join(output_folder, 'imsmanifest.xml')
        with open(manifest_path, 'w') as f:
            f.write(scorm_manifest)

        scorm_zip = os.path.join(output_folder, 'scorm_package.zip')
        with zipfile.ZipFile(scorm_zip, 'w') as zipf:
            for image in slide_images:
                zipf.write(image, os.path.basename(image))
            zipf.write(manifest_path, 'imsmanifest.xml')
            for i in range(len(slide_images)):
                html_file = os.path.join(output_folder, f"slide_{i + 1}.html")
                zipf.write(html_file, os.path.basename(html_file))

        print(f"SCORM package created at {scorm_zip}")
    except Exception as e:
        print(f"Error converting PPTX to SCORM: {e}")
        raise

def generate_html_files(slide_images, output_folder, title, author, slide_titles, use_mini_images, style="black", fullscreen=False):
    try:
        for i, image in enumerate(slide_images):
            if use_mini_images:
                thumbnails = "".join([
                    f'<div style="text-align: center; margin-bottom: 10px;"><a href="slide_{j + 1}.html"><img src="{os.path.basename(img)}" alt="Slide {j + 1}" style="width:100%; background-color: #ccc; {"border: 6px solid red; box-shadow: 0 0 10px 2px rgba(255, 0, 0, 0.5);" if i == j else "border: 2px solid transparent;"}" id="thumb_{j + 1}"></a><br><span style="color: {"white" if style in ["black", "black-red", "green"] else "black"};">{slide_titles[j] if slide_titles else j + 1}</span></div>'
                    for j, img in enumerate(slide_images)])
            else:
                thumbnails = "".join([
                    f'<div style="text-align: center; margin-bottom: 10px;"><a href="slide_{j + 1}.html"><span style="color: {"white" if style in ["black", "black-red", "green"] else "black"};" id="thumb_{j + 1}">{slide_titles[j] if slide_titles else j + 1}</span></a></div>'
                    for j, img in enumerate(slide_images)])
            next_slide = f"slide_{i + 2}.html" if i + 1 < len(slide_images) else "slide_1.html"
            prev_slide = f"slide_{i}.html" if i > 0 else f"slide_{len(slide_images)}.html"
            if style == "black":
                background_color = "#1a1a1a"
                text_color = "white"
                title_color = "white"
            elif style == "white":
                background_color = "white"
                text_color = "black"
                title_color = "black"
            elif style == "green":
                background_color = "#013220"  # Bottle green color
                text_color = "white"
                title_color = "white"  # Ensuring the title is visible in green style
            elif style == "black-red":
                background_color = "linear-gradient(to bottom, black, red)"
                text_color = "white"
                title_color = "white"

            sidebar_display = "none" if fullscreen else "block"

            html_content = f'''
            <!DOCTYPE html>
            <html>
            <head>
                <title>{title}</title>
                <style>
                    body {{
                        display: flex;
                        height: 100vh;
                        margin: 0;
                        background: {background_color};
                        color: {text_color};
                        font-family: Arial, sans-serif;
                    }}
                    .content {{
                        display: flex;
                        justify-content: center;
                        align-items: center;
                        flex-direction: column;
                        flex-grow: 1;
                        position: relative;
                    }}
                    img {{
                        max-width: 90%;
                        max-height: 70vh;
                    }}
                    .nav-buttons {{
                        margin-top: 20px;
                    }}
                    .nav-buttons button {{
                        padding: 6px 15px;  /* significantly reduced */
                        font-size: 14px;  /* slightly reduced */
                        margin: 0 6px;  /* significantly reduced */
                        cursor: pointer;
                        background-color: #444;
                        color: white;
                        border: none;
                        border-radius: 5px;
                        transition: background-color 0.3s ease;
                    }}
                    .nav-buttons button:hover {{
                        background-color: #1abc9c;
                    }}
                    .sidebar {{
                        width: 200px;
                        background-color: #333;
                        padding: 10px;
                        overflow-y: auto;
                        height: 100vh;
                        display: {sidebar_display};
                    }}
                    .sidebar img {{
                        transition: border 0.3s ease, background-color 0.3s ease;
                        background-color: #ccc;
                    }}
                    .sidebar img:hover {{
                        border: 2px solid #1abc9c;
                    }}
                    .sidebar img.active {{
                        border: 6px solid red;
                        box-shadow: 0 0 10px 2px rgba(255, 0, 0, 0.5);
                        background-color: #666;
                    }}
                    .sidebar span {{
                        transition: border 0.3s ease, background-color 0.3s ease;
                        background-color: transparent;
                        display: block;
                        text-align: center;
                        padding: 10px;
                        font-size: 14px;  /* slightly reduced */
                        color: {"white" if style == "green" else text_color};  /* Ensure white color for green style */
                    }}
                    .sidebar span:hover {{
                        background-color: transparent;
                    }}
                    .sidebar span.active {{
                        background-color: transparent;
                    }}
                    .slide-number {{
                        position: absolute;
                        bottom: 20px;
                        right: 20px;
                        background-color: rgba(0, 0, 0, 0.5);
                        color: white;
                        padding: 5px 10px;
                        border-radius: 5px;
                        font-size: 16px;
                    }}
                    .slide-title {{
                        position: absolute;
                        top: 10px;
                        left: 50%;
                        transform: translateX(-50%);
                        background-color: rgba(0, 0, 0, 0.5);
                        color: {title_color};
                        padding: 5px 10px;
                        border-radius: 5px;
                        font-size: 32px;  /* increased */
                        text-align: center;
                        width: 80%;
                    }}
                    .author-note {{
                        position: absolute;
                        bottom: 20px;
                        left: 20px;
                        background-color: rgba(0, 0, 0, 0.5);
                        color: white;
                        padding: 5px 10px;
                        border-radius: 5px;
                        font-size: 16px;
                    }}
                </style>
                <script>
                    window.onload = function() {{
                        var currentThumb = document.getElementById("thumb_{i + 1}");
                        if (currentThumb) {{
                            currentThumb.scrollIntoView({{behavior: "smooth", block: "center"}});
                            currentThumb.classList.add("active");
                        }}
                    }};
                </script>
            </head>
            <body>
                <div class="sidebar">
                    {thumbnails}
                </div>
                <div class="content">
                    <div class="slide-title">{title}</div>
                    <img src="{os.path.basename(image)}" alt="Slide {i + 1}">
                    <div class="slide-number">Slide {i + 1} of {len(slide_images)}</div>
                    <div class="nav-buttons">
                        <button onclick="location.href='{prev_slide}'">Wstecz</button>
                        <button onclick="location.href='{next_slide}'">Dalej</button>
                    </div>
                    <div class="author-note">Autor prezentacji: {author}</div>
                </div>
            </body>
            </html>
            '''
            html_path = os.path.join(output_folder, f"slide_{i + 1}.html")
            with open(html_path, 'w') as f:
                f.write(html_content)
    except Exception as e:
        print(f"Error generating HTML files: {e}")
        raise

def generate_scorm_manifest(slide_images):
    try:
        manifest_template = '''<?xml version="1.0" encoding="UTF-8"?>
    <manifest identifier="com.example.scorm" version="1.2">
      <metadata>
        <schema>ADL SCORM</schema>
        <schemaversion>1.2</schemaversion>
      </metadata>
      <organizations default="ORG">
        <organization identifier="ORG">
          <title>Sample SCORM Course</title>
          {items}
        </organization>
      </organizations>
      <resources>
        {resources}
      </resources>
    </manifest>'''

        item_template = '<item identifier="ITEM-{index}" identifierref="RES-{index}"><title>Slide {index}</title></item>'
        resource_template = '<resource identifier="RES-{index}" type="webcontent" href="slide_{index}.html"><file href="slide_{index}.html"/><file href="slide_{index}.png"/></resource>'

        items = '\n'.join([item_template.format(index=i + 1) for i in range(len(slide_images))])
        resources = '\n'.join([resource_template.format(index=i + 1) for i in range(len(slide_images))])

        return manifest_template.format(items=items, resources=resources)
    except Exception as e:
        print(f"Error generating SCORM manifest: {e}")
        raise

class SCORMConverterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("KONWERTER SCROM")
        self.root.geometry("1170x780")  # 30% increase from 900x600
        self.root.resizable(False, False)

        self.ppsx_file = ""
        self.pptx_file = ""
        self.output_folder = ""
        self.style = "black"
        self.fullscreen = tk.BooleanVar()
        self.use_slide_titles = tk.BooleanVar()
        self.use_mini_images = tk.BooleanVar()

        # Set background image
        self.set_background_image()

        # Set custom fonts
        self.title_font = font.Font(family="Helvetica", size=24, weight="bold")
        self.button_font = font.Font(family="Helvetica", size=12, weight="bold")  # slightly reduced
        self.label_font = font.Font(family="Helvetica", size=10)  # reduced
        self.entry_font = font.Font(family="Helvetica", size=10)  # reduced

        self.create_widgets()

    def set_background_image(self):
        background_image_path = resource_path("background.jpg")
        self.background_image = Image.open(background_image_path)
        self.background_image = self.background_image.resize((1170, 780), Image.LANCZOS)  # 30% increase from 900x600
        self.background_photo = ImageTk.PhotoImage(self.background_image)

        self.background_label = tk.Label(self.root, image=self.background_photo)
        self.background_label.place(relx=0.5, rely=0.5, anchor='center')

    def create_widgets(self):
        title_frame = tk.Frame(self.root, bg="lightblue", bd=5, relief="ridge")
        title_frame.pack(pady=10, padx=10, fill="x")

        title_label = tk.Label(title_frame, text="KONWERTER SCROM", font=self.title_font, bg="lightblue")
        title_label.pack(padx=10, pady=10)

        button_style = {"font": self.button_font, "bg": "lightgrey", "relief": "raised", "bd": 4, "width": 18, "height": 1}  # slightly reduced

        button_frame = tk.Frame(self.root, bg="lightblue")
        button_frame.pack(side=tk.LEFT, padx=20, pady=20, fill=tk.Y)

        style = ttk.Style()
        style.configure("TButton", padding=6, relief="flat", background="#ccc", foreground="#000",
                        font=("Helvetica", 12, "bold"), borderwidth=0)
        style.map("TButton", background=[("active", "#aaa")])

        ttk.Button(button_frame, text="Wybierz plik", command=self.select_ppsx_file).pack(pady=10)
        self.ppsx_label = tk.Label(button_frame, text="Nie wybrano pliku", font=self.label_font, bg="lightblue")
        self.ppsx_label.pack(pady=10)

        ttk.Button(button_frame, text="Wybierz folder", command=self.select_output_folder).pack(pady=10)
        self.output_label = tk.Label(button_frame, text="Nie wybrano folderu", font=self.label_font, bg="lightblue")
        self.output_label.pack(pady=10)

        style_frame = tk.Frame(button_frame, bg="lightblue")
        style_frame.pack(pady=10)

        # Customize style selection button
        tk.Label(style_frame, text="Wybierz Styl:", font=self.label_font, bg="lightblue").pack(side=tk.LEFT, padx=5)

        self.style_combobox = ttk.Combobox(style_frame, values=["black", "white", "green", "black-red"], font=self.label_font)
        self.style_combobox.set("black")
        self.style_combobox.pack(side=tk.LEFT, padx=5)

        # Style the combobox dropdown
        style = ttk.Style()
        style.theme_use('clam')
        style.configure("TCombobox", fieldbackground="white", background="lightgrey", arrowcolor="black", bordercolor="lightgrey")
        style.map('TCombobox', fieldbackground=[('readonly', 'lightgrey')],
                  background=[('readonly', 'lightgrey')],
                  arrowcolor=[('active', 'black')],
                  bordercolor=[('active', 'lightgrey')])

        options_frame = tk.Frame(button_frame, bg="lightblue")
        options_frame.pack(pady=10, fill=tk.X)

        tk.Checkbutton(options_frame, text="Tryb pełnoekranowy", variable=self.fullscreen, font=self.label_font,
                       bg="lightblue").grid(row=0, column=0, sticky='w', pady=5)

        tk.Checkbutton(options_frame, text="Użyj tytułów slajdów", variable=self.use_slide_titles, font=self.label_font,
                       bg="lightblue").grid(row=1, column=0, sticky='w', pady=5)

        tk.Checkbutton(options_frame, text="Użyj mini zdjęć w menu", variable=self.use_mini_images, font=self.label_font,
                       bg="lightblue").grid(row=2, column=0, sticky='w', pady=5)

        self.title_label = tk.Label(button_frame, text="Tytuł prezentacji:", font=self.label_font, bg="lightblue")
        self.title_label.pack(pady=10)
        self.title_entry = tk.Entry(button_frame, font=self.entry_font, width=25)
        self.title_entry.pack(pady=10)

        self.author_label = tk.Label(button_frame, text="Autor prezentacji:", font=self.label_font, bg="lightblue")
        self.author_label.pack(pady=10)
        self.author_entry = tk.Entry(button_frame, font=self.entry_font, width=25)
        self.author_entry.pack(pady=10)

        ttk.Button(button_frame, text="Konwertuj", command=self.convert).pack(pady=20)

        # Add version label at the bottom right corner
        self.footer_label = tk.Label(self.root, text="ver. 1.1", font=self.label_font, bg="lightblue")
        self.footer_label.place(relx=0.0, rely=1.0, anchor='sw', x=10, y=-10)

        # Add label at the bottom left corner
        self.version_label = tk.Label(self.root, text="Norbert Kowalski CUI Wrocław", font=self.label_font,
                                     bg="lightblue")
        self.version_label.place(relx=1.0, rely=1.0, anchor='se', x=-10, y=-10)

    def select_ppsx_file(self):
        self.ppsx_file = filedialog.askopenfilename(filetypes=[("PPSX files", "*.ppsx")])
        self.pptx_file = self.ppsx_file.replace(".ppsx", ".pptx")
        self.ppsx_label.config(text=os.path.basename(self.ppsx_file))

    def select_output_folder(self):
        self.output_folder = filedialog.askdirectory()
        self.output_label.config(text=self.output_folder)

    def convert(self):
        if not self.ppsx_file or not self.output_folder:
            messagebox.showerror("Błąd", "Proszę wybrać plik PPSX i folder wyjściowy.")
            return

        self.style = self.style_combobox.get()
        self.fullscreen_mode = self.fullscreen.get()
        self.title = self.title_entry.get()
        self.author = self.author_entry.get()
        self.use_slide_titles_mode = self.use_slide_titles.get()
        self.use_mini_images_mode = self.use_mini_images.get()

        try:
            # Convert PPSX to PPTX
            convert_ppsx_to_pptx(self.ppsx_file, self.pptx_file)
            # Convert PPTX to SCORM with the selected style, title, and author
            convert_pptx_to_scorm(self.pptx_file, self.output_folder, self.title, self.author, self.use_slide_titles_mode, self.use_mini_images_mode, self.style, self.fullscreen_mode)
            messagebox.showinfo("Sukces", "Konwersja zakończona pomyślnie.")
        except Exception as e:
            messagebox.showerror("Błąd", f"Wystąpił błąd: {e}")


if __name__ == "__main__":
    root = tk.Tk()
    app = SCORMConverterApp(root)
    root.mainloop()
