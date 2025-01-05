import os
import zipfile
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import win32com.client


def convert_ppsx_to_pptx(ppsx_file, pptx_file):
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(ppsx_file)
    presentation.SaveAs(pptx_file, 24)  # 24 oznacza format pptx
    presentation.Close()
    powerpoint.Quit()


def save_slide_as_image(slide, path, width, height):
    image = Image.new('RGB', (width, height), color='white')
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left = int(shape.left * 96 / 914400)  # EMU to pixels
            top = int(shape.top * 96 / 914400)  # EMU to pixels
            with io.BytesIO(shape.image.blob) as image_stream:
                pic = Image.open(image_stream)
                image.paste(pic, (left, top))
    image.save(path, 'PNG')


def convert_pptx_to_images(pptx_file, output_folder):
    prs = Presentation(pptx_file)

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    slide_images = []
    for i, slide in enumerate(prs.slides):
        image_path = os.path.join(output_folder, f"slide_{i + 1}.png")
        width = prs.slide_width // 914400 * 96  # EMU to pixels
        height = prs.slide_height // 914400 * 96  # EMU to pixels
        save_slide_as_image(slide, image_path, width, height)
        slide_images.append(image_path)
    return slide_images


def convert_pptx_to_scorm(pptx_file, output_folder):
    slide_images = convert_pptx_to_images(pptx_file, output_folder)

    scorm_manifest = generate_scorm_manifest(slide_images)
    manifest_path = os.path.join(output_folder, 'imsmanifest.xml')
    with open(manifest_path, 'w') as f:
        f.write(scorm_manifest)

    scorm_zip = os.path.join(output_folder, 'scorm_package.zip')
    with zipfile.ZipFile(scorm_zip, 'w') as zipf:
        for image in slide_images:
            zipf.write(image, os.path.basename(image))
        zipf.write(manifest_path, 'imsmanifest.xml')

    print(f"SCORM package created at {scorm_zip}")


def generate_scorm_manifest(slide_images):
    manifest_template = '''<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="com.example.scorm" version="1.2">
  <metadata>
    <schema>ADL SCORM</schema>
    <schemaversion>1.2</schemaversion>
  </metadata>
  <organizations default="ORG">
    <organization identifier="ORG">
      <title>Sample SCORM Course</title>
      <item identifier="ITEM-1" identifierref="RES-1">
        <title>Sample SCORM Course</title>
      </item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="RES-1" type="webcontent" href="slide_1.png">
      {resources}
    </resource>
  </resources>
</manifest>'''

    resource_template = '<file href="{filename}"/>'
    resources = '\n'.join([resource_template.format(filename=os.path.basename(image)) for image in slide_images])
    return manifest_template.format(resources=resources)


# Ścieżka do pliku PPSX
ppsx_file = r'C:\Users\kowal\OneDrive\Pulpit\CUI Prezentacja\eLearning_CUI_WrocławEDUEOFVer.ppsx'
# Ścieżka do zapisu pliku PPTX
pptx_file = r'C:\Users\kowal\OneDrive\Pulpit\CUI Prezentacja\eLearning_CUI_WrocławEDUEOFVer.pptx'

convert_ppsx_to_pptx(ppsx_file, pptx_file)

output_folder = 'output_scorm'
convert_pptx_to_scorm(pptx_file, output_folder)
