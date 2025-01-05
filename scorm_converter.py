import os
import zipfile
import io
import time
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import win32com.client
from tkinter import Tk, filedialog, messagebox, Label, Button, Canvas
from PIL import Image, ImageTk


def close_powerpoint_instances():
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Quit()
        print("Closed existing PowerPoint instances.")
    except Exception as e:
        print(f"Error closing PowerPoint instances: {e}")


def convert_ppsx_to_pptx(ppsx_file, pptx_file):
    close_powerpoint_instances()
    time.sleep(5)  # Give more time to ensure PowerPoint instances are closed
    try:
        if not os.path.exists(ppsx_file):
            raise FileNotFoundError(f"File {ppsx_file} not found.")
        if not os.access(ppsx_file, os.R_OK):
            raise PermissionError(f"File {ppsx_file} is not readable.")

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        if powerpoint is None:
            raise Exception("Failed to start PowerPoint.")

        powerpoint.Visible = 1
        print(f"Opening PPSX file: {ppsx_file}")
        presentation = powerpoint.Presentations.Open(ppsx_file, WithWindow=False)
        print(f"Saving as PPTX file: {pptx_file}")
        presentation.SaveAs(pptx_file, 24)  # 24 oznacza format pptx
        presentation.Close()
        powerpoint.Quit()
        print(f"Converted {ppsx_file} to {pptx_file}")
    except Exception as e:
        messagebox.showerror("Błąd", f"Błąd podczas konwersji PPSX na PPTX: {e}")
        print(f"Error converting PPSX to PPTX: {e}")


def save_slide_as_image(slide, path, width, height):
    image = Image.new('RGB', (width, height), color='white')
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left = int(shape.left * 96 / 914400)  # EMU to pixels
            top = int(shape.top * 96 / 914400)  # EMU to pixels
            with io.BytesIO(shape.image.blob) as image_stream:
                pic = Image.open(image_stream)
                image.paste(pic, (left, top))
    image.save(path, 'PNG')


def convert_pptx_to_images(pptx_file, output_folder):
    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        raise Exception(f"Error loading PPTX file: {e}")

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    slide_images = []
    for i, slide in enumerate(prs.slides):
        image_path = os.path.join(output_folder, f"slide_{i + 1}.png")
        width = prs.slide_width // 914400 * 96  # EMU to pixels
        height = prs.slide_height // 914400 * 96  # EMU to pixels
        save_slide_as_image(slide, image_path, width, height)
        slide_images.append(image_path)
    return slide_images


def convert_pptx_to_scorm(pptx_file, output_folder):
    try:
        slide_images = convert_pptx_to_images(pptx_file, output_folder)

        scorm_manifest = generate_scorm_manifest(slide_images)
        manifest_path = os.path.join(output_folder, 'imsmanifest.xml')
        with open(manifest_path, 'w') as f:
            f.write(scorm_manifest)

        scorm_zip = os.path.join(output_folder, 'scorm_package.zip')
        with zipfile.ZipFile(scorm_zip, 'w') as zipf:
            for image in slide_images:
                zipf.write(image, os.path.basename(image))
            zipf.write(manifest_path, 'imsmanifest.xml')

        print(f"SCORM package created at {scorm_zip}")
        messagebox.showinfo("Sukces", f"Pakiet SCORM został utworzony w {scorm_zip}")
    except Exception as e:
        messagebox.showerror("Błąd", f"Błąd podczas konwersji PPTX na SCORM: {e}")
        print(f"Error converting PPTX to SCORM: {e}")


def generate_scorm_manifest(slide_images):
    manifest_template = '''<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="com.example.scorm" version="1.2">
  <metadata>
    <schema>ADL SCORM</schema>
    <schemaversion>1.2</schemaversion>
  </metadata>
  <organizations default="ORG">
    <organization identifier="ORG">
      <title>Sample SCORM Course</title>
      <item identifier="ITEM-1" identifierref="RES-1">
        <title>Sample SCORM Course</title>
      </item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="RES-1" type="webcontent" href="slide_1.png">
      {resources}
    </resource>
  </resources>
</manifest>'''

    resource_template = '<file href="{filename}"/>'
    resources = '\n'.join([resource_template.format(filename=os.path.basename(image)) for image in slide_images])
    return manifest_template.format(resources=resources)


def select_ppsx_file():
    ppsx_file = filedialog.askopenfilename(title="Wybierz plik PPSX", filetypes=[("PPSX files", "*.ppsx")])
    if ppsx_file:
        ppsx_label.config(text=ppsx_file)
    return ppsx_file


def select_pptx_file():
    pptx_file = filedialog.askopenfilename(title="Wybierz plik PPTX", filetypes=[("PPTX files", "*.pptx")])
    if pptx_file:
        pptx_label.config(text=pptx_file)
    return pptx_file


def select_output_folder():
    output_folder = filedialog.askdirectory(title="Wybierz folder wyjściowy")
    if output_folder:
        output_label.config(text=output_folder)
    return output_folder


def start_conversion():
    ppsx_file = ppsx_label.cget("text")
    pptx_file = pptx_label.cget("text")
    output_folder = output_label.cget("text")

    if not ppsx_file and not pptx_file:
        messagebox.showerror("Błąd", "Nie wybrano pliku PPSX ani PPTX.")
        return

    if not output_folder:
        messagebox.showerror("Błąd", "Nie wybrano folderu wyjściowego.")
        return

    if ppsx_file:
        pptx_file = ppsx_file.replace(".ppsx", ".pptx")
        convert_ppsx_to_pptx(ppsx_file, pptx_file)

    convert_pptx_to_scorm(pptx_file, output_folder)


def main():
    root = Tk()
    root.title("PPSX to SCORM Converter")
    root.geometry("700x600")

    # Load background image
    bg_image_path = "background.jpg"
    if not os.path.exists(bg_image_path):
        messagebox.showerror("Błąd", "Plik background.jpg nie został znaleziony.")
        return

    bg_image = Image.open(bg_image_path)
    bg_image = bg_image.resize((700, 600), Image.LANCZOS)
    bg_photo = ImageTk.PhotoImage(bg_image)

    # Create a canvas
    canvas = Canvas(root, width=700, height=600)
    canvas.pack(fill="both", expand=True)
    canvas.create_image(0, 0, image=bg_photo, anchor="nw")

    # Add labels and buttons
    canvas.create_text(350, 50, text="PPSX to SCORM Converter", font=("Helvetica", 24), fill="white")
    canvas.create_text(150, 150, text="Wybierz plik PPSX:", font=("Helvetica", 14), fill="white")
    canvas.create_text(150, 200, text="Wybierz plik PPTX:", font=("Helvetica", 14), fill="white")
    canvas.create_text(150, 300, text="Wybierz folder do zapisu:", font=("Helvetica", 14), fill="white")
    canvas.create_text(350, 550, text="Auto program Norbert Kowalski CUI Wrocław", font=("Helvetica", 10), fill="white")

    global ppsx_label, pptx_label, output_label
    ppsx_label = Label(root, text="", bg="#1a1a1a", fg="white", font=("Helvetica", 12))
    pptx_label = Label(root, text="", bg="#1a1a1a", fg="white", font=("Helvetica", 12))
    output_label = Label(root, text="", bg="#1a1a1a", fg="white", font=("Helvetica", 12))

    ppsx_button = Button(root, text="Wybierz plik PPSX", command=select_ppsx_file)
    pptx_button = Button(root, text="Wybierz plik PPTX", command=select_pptx_file)
    output_button = Button(root, text="Wybierz folder", command=select_output_folder)
    convert_button = Button(root, text="Konwertuj", command=start_conversion)

    canvas.create_window(50, 150, anchor="nw", window=ppsx_button)
    canvas.create_window(50, 200, anchor="nw", window=pptx_button)
    canvas.create_window(50, 300, anchor="nw", window=output_button)
    canvas.create_window(150, 400, anchor="nw", window=convert_button)
    canvas.create_window(250, 150, anchor="nw", window=ppsx_label)
    canvas.create_window(250, 200, anchor="nw", window=pptx_label)
    canvas.create_window(250, 300, anchor="nw", window=output_label)

    root.mainloop()


if __name__ == "__main__":
    main()
